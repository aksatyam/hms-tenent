#!/usr/bin/env python3
"""
HMS Enterprise — Client Pitch Deck (Enterprise Grade)
Rich-themed PPTX generator with strong TechDigital WishTree branding.
Based on HMS_Enterprise_Product_Document.docx
Author: Ashish Kumar Satyam | Organization: TechDigital WishTree
Date: 16 March 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ─── Output ───
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BASE_DIR = os.path.dirname(SCRIPT_DIR)
OUTPUT_DIR = os.path.join(BASE_DIR, "docs", "presentations")
os.makedirs(OUTPUT_DIR, exist_ok=True)
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "HMS_Enterprise_Client_Pitch_Deck.pptx")

# ─── Branding ───
COMPANY_FULL = "TechDigital WishTree"
COMPANY_SHORT = "TGWT"
PRODUCT_NAME = "HMS Enterprise"
AUTHOR = "Ashish Kumar Satyam"
DOC_DATE = "16 March 2026"

# ─── Brand Colors ───
NAVY       = RGBColor(27, 58, 92)
DARK_NAVY  = RGBColor(15, 38, 64)
TEAL       = RGBColor(13, 115, 119)
TEAL_LIGHT = RGBColor(16, 185, 129)
GOLD       = RGBColor(196, 154, 42)
AI_PURPLE  = RGBColor(124, 58, 237)
AI_BLUE    = RGBColor(79, 70, 229)
WHITE      = RGBColor(255, 255, 255)
BLACK      = RGBColor(26, 32, 44)
GRAY       = RGBColor(113, 128, 150)
LIGHT_GRAY = RGBColor(237, 242, 247)
RED        = RGBColor(229, 62, 62)
ORANGE     = RGBColor(237, 137, 54)
GREEN      = RGBColor(56, 161, 105)
BODY_TEXT  = RGBColor(45, 55, 72)
LIGHT_PURPLE = RGBColor(167, 139, 250)
SOFT_GRAY  = RGBColor(160, 174, 192)

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, bold=False, italic=False, color=WHITE, font='Calibri', align=PP_ALIGN.LEFT, line_spacing=1.2):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(size * line_spacing)
    return p

def add_para(tf, text, size=16, bold=False, italic=False, color=WHITE, font='Calibri', align=PP_ALIGN.LEFT, space_before=0, space_after=4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_bullet_para(tf, text, size=14, color=WHITE, bold_prefix=None, font='Calibri Light', space_after=6, indent_level=0):
    p = tf.add_paragraph()
    p.level = indent_level
    if bold_prefix:
        run1 = p.add_run()
        run1.text = bold_prefix + " "
        run1.font.size = Pt(size)
        run1.font.bold = True
        run1.font.color.rgb = color
        run1.font.name = 'Calibri'
        run2 = p.add_run()
        run2.text = text
        run2.font.size = Pt(size)
        run2.font.bold = False
        run2.font.color.rgb = color
        run2.font.name = font
    else:
        run = p.add_run()
        run.text = "\u2022  " + text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    p.space_after = Pt(space_after)
    p.space_before = Pt(2)
    return p

def add_slide_footer(slide):
    """Add branded footer bar."""
    # Footer bar
    add_shape(slide, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), DARK_NAVY)
    # Gold accent line
    add_shape(slide, Inches(0), Inches(7.02), SLIDE_W, Inches(0.03), GOLD)
    # Company name
    tb = add_text_box(slide, Inches(0.5), Inches(7.08), Inches(5), Inches(0.35))
    set_text(tb.text_frame, f"{COMPANY_FULL}  |  {PRODUCT_NAME}", size=9, color=SOFT_GRAY, font='Calibri Light')
    # Confidential
    tb2 = add_text_box(slide, Inches(9), Inches(7.08), Inches(4), Inches(0.35))
    set_text(tb2.text_frame, "CONFIDENTIAL", size=9, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

def add_slide_header(slide, title, subtitle=None, icon=""):
    """Standard navy header bar with title."""
    # Header bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_NAVY)
    # Gold accent
    add_shape(slide, Inches(0), Inches(1.15), SLIDE_W, Inches(0.05), GOLD)
    # Title
    tb = add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7))
    tf = tb.text_frame
    if icon:
        set_text(tf, f"{icon}  {title}", size=28, bold=True, color=WHITE)
    else:
        set_text(tf, title, size=28, bold=True, color=WHITE)
    # Subtitle
    if subtitle:
        tb2 = add_text_box(slide, Inches(0.6), Inches(0.75), Inches(10), Inches(0.35))
        set_text(tb2.text_frame, subtitle, size=12, color=SOFT_GRAY, font='Calibri Light')
    # TGWT logo text
    tb3 = add_text_box(slide, Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5))
    set_text(tb3.text_frame, COMPANY_SHORT, size=16, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

def add_stat_box(slide, left, top, value, label, width=Inches(2.5), bg_color=DARK_NAVY, val_color=LIGHT_PURPLE):
    """Add a stat number box."""
    box = add_shape(slide, left, top, width, Inches(1.1), bg_color)
    box.shadow.inherit = False
    # Value
    tb_val = add_text_box(slide, left, top + Inches(0.05), width, Inches(0.6))
    set_text(tb_val.text_frame, value, size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    # Label
    tb_lbl = add_text_box(slide, left, top + Inches(0.6), width, Inches(0.4))
    set_text(tb_lbl.text_frame, label.upper(), size=8, bold=True, color=SOFT_GRAY, align=PP_ALIGN.CENTER, font='Calibri')
    return box


# ═══════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title Slide — Full navy background with branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full navy background
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK_NAVY)

    # Top teal accent line
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), TEAL)

    # Company name
    tb = add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.6))
    set_text(tb.text_frame, COMPANY_FULL.upper(), size=18, bold=True, color=GOLD, font='Calibri', align=PP_ALIGN.CENTER)

    # Tagline
    tb2 = add_text_box(slide, Inches(1), Inches(1.7), Inches(11), Inches(0.4))
    set_text(tb2.text_frame, "Technology  \u2022  Digital Transformation  \u2022  Healthcare Innovation", size=11, color=SOFT_GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)

    # Gold divider
    add_shape(slide, Inches(4), Inches(2.3), Inches(5.3), Inches(0.04), GOLD)

    # Product name
    tb3 = add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.0))
    set_text(tb3.text_frame, "HMS ENTERPRISE", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    tb4 = add_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.6))
    set_text(tb4.text_frame, "Client Pitch Deck", size=26, color=LIGHT_PURPLE, font='Calibri Light', align=PP_ALIGN.CENTER)

    # Description
    tb5 = add_text_box(slide, Inches(2), Inches(4.5), Inches(9), Inches(0.5))
    set_text(tb5.text_frame, "Cloud-Native Multi-Tenant Hospital Management Platform", size=14, color=SOFT_GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)

    tb6 = add_text_box(slide, Inches(2), Inches(4.9), Inches(9), Inches(0.4))
    set_text(tb6.text_frame, "12 Integrated Modules  \u2022  AI/ML Innovation  \u2022  ABDM & DPDP Compliant", size=11, color=GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)

    # Gold divider bottom
    add_shape(slide, Inches(4), Inches(5.6), Inches(5.3), Inches(0.04), GOLD)

    # Author & Date
    tb7 = add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.4))
    set_text(tb7.text_frame, f"Presented by: {AUTHOR}", size=12, color=WHITE, font='Calibri Light', align=PP_ALIGN.CENTER)

    tb8 = add_text_box(slide, Inches(1), Inches(6.35), Inches(11), Inches(0.4))
    set_text(tb8.text_frame, f"{COMPANY_FULL}  |  {DOC_DATE}", size=10, color=GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)

    # Bottom teal accent
    add_shape(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), TEAL)


def slide_02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Agenda", "What We'll Cover Today")
    add_slide_footer(slide)

    items = [
        ("01", "The Problem", "Healthcare IT challenges in India"),
        ("02", "Market Opportunity", "\u20B970,000 Cr+ market with massive digital gap"),
        ("03", "Our Solution", "HMS Enterprise platform overview"),
        ("04", "Platform Architecture", "Cloud-native, multi-tenant design"),
        ("05", "12 Integrated Modules", "Complete hospital workflow coverage"),
        ("06", "AI/ML Innovation", "12 AI modules across 4 phases"),
        ("07", "Security & Compliance", "DPDP, ABDM, NABH built-in"),
        ("08", "Go-to-Market & Pricing", "3-tier SaaS model"),
        ("09", "Implementation Roadmap", "54-week delivery timeline"),
        ("10", "Why TechDigital WishTree", "Our competitive edge"),
    ]

    y = Inches(1.5)
    for num, title, desc in items:
        # Number box
        add_shape(slide, Inches(0.8), y, Inches(0.7), Inches(0.45), DARK_NAVY)
        tb_num = add_text_box(slide, Inches(0.8), y + Inches(0.02), Inches(0.7), Inches(0.4))
        set_text(tb_num.text_frame, num, size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        # Title
        tb_title = add_text_box(slide, Inches(1.7), y, Inches(4), Inches(0.45))
        set_text(tb_title.text_frame, title, size=14, bold=True, color=NAVY)
        # Desc
        tb_desc = add_text_box(slide, Inches(5.5), y, Inches(6), Inches(0.45))
        set_text(tb_desc.text_frame, desc, size=12, color=GRAY, font='Calibri Light')
        y += Inches(0.52)


def slide_03_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "The Problem", "Healthcare IT in India is broken")
    add_slide_footer(slide)

    problems = [
        ("\u26A0\uFE0F", "Legacy Systems", "85% of hospitals run on-premise, single-tenant software from the 2000s. No cloud, no interoperability, no mobile access."),
        ("\U0001F512", "Data Silos", "Patient records scattered across departments. No unified view. No clinical decision support. Paper-based referrals."),
        ("\U0001F4B8", "Revenue Leakage", "15\u201320% revenue loss due to manual billing, coding errors, missed charges, and claim denials. No AI-powered RCM."),
        ("\u274C", "No ABDM Compliance", "ABDM mandatory by 2026 but most vendors have zero FHIR/ABHA integration. Hospitals face regulatory risk."),
        ("\U0001F916", "Zero AI Adoption", "No ambient clinical AI, no CDSS, no predictive analytics. Doctors spend 40% of time on documentation."),
        ("\U0001F6E1\uFE0F", "Security Gaps", "DPDP Act 2023 compliance missing. No field-level encryption. No consent management. Audit trails incomplete."),
    ]

    y = Inches(1.5)
    for i, (icon, title, desc) in enumerate(problems):
        col = Inches(0.5) if i % 2 == 0 else Inches(6.8)
        # Card background
        card = add_shape(slide, col, y, Inches(5.8), Inches(1.0), RGBColor(247, 250, 252))
        # Left accent
        add_shape(slide, col, y, Inches(0.06), Inches(1.0), RED if i < 2 else ORANGE if i < 4 else TEAL)
        # Icon + Title
        tb = add_text_box(slide, col + Inches(0.2), y + Inches(0.08), Inches(5.3), Inches(0.35))
        set_text(tb.text_frame, f"{icon}  {title}", size=13, bold=True, color=NAVY)
        # Description
        tb2 = add_text_box(slide, col + Inches(0.2), y + Inches(0.42), Inches(5.3), Inches(0.55))
        set_text(tb2.text_frame, desc, size=10, color=BODY_TEXT, font='Calibri Light')
        if i % 2 == 1:
            y += Inches(1.15)


def slide_04_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Market Opportunity", "India's healthcare digitization is at an inflection point")
    add_slide_footer(slide)

    # Stat boxes
    stats = [
        ("\u20B970,000Cr+", "Market Size 2030"),
        ("18.7%", "CAGR Growth"),
        ("70,000+", "Hospitals in India"),
        ("15%", "Digital Adoption"),
    ]
    x = Inches(0.7)
    for val, label in stats:
        add_stat_box(slide, x, Inches(1.5), val, label, Inches(2.8))
        x += Inches(3.1)

    # Key insights
    tb = add_text_box(slide, Inches(0.7), Inches(3.0), Inches(12), Inches(0.5))
    set_text(tb.text_frame, "Key Market Insights", size=18, bold=True, color=NAVY)

    insights_left = [
        "Only 2\u20133 Indian-origin vendors offer true multi-tenant SaaS",
        "Zero vendors provide ambient clinical AI or CDSS",
        "ABDM compliance becomes mandatory in 2026 — most vendors are unprepared",
        "Mid-market segment (100\u2013500 beds) is vastly underserved",
    ]
    insights_right = [
        "No Indian vendor offers AI-powered revenue cycle management",
        "Predictive bed management is a blue-sky opportunity",
        "Government digitization (NHM, PMJAY) creating massive demand",
        "18 validated market gaps — 4 classified as critical",
    ]

    tb_l = add_text_box(slide, Inches(0.7), Inches(3.5), Inches(5.8), Inches(3.2))
    tf = tb_l.text_frame
    tf.word_wrap = True
    set_text(tf, "", size=1)
    for item in insights_left:
        add_bullet_para(tf, item, size=11, color=BODY_TEXT, font='Calibri Light')

    tb_r = add_text_box(slide, Inches(6.8), Inches(3.5), Inches(5.8), Inches(3.2))
    tf2 = tb_r.text_frame
    tf2.word_wrap = True
    set_text(tf2, "", size=1)
    for item in insights_right:
        add_bullet_para(tf2, item, size=11, color=BODY_TEXT, font='Calibri Light')

    # Callout
    callout = add_shape(slide, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.7), RGBColor(250, 245, 255))
    add_shape(slide, Inches(0.7), Inches(5.8), Inches(0.06), Inches(0.7), AI_PURPLE)
    tb_c = add_text_box(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.6))
    set_text(tb_c.text_frame, "\u26A1  HMS Enterprise addresses all 18 market gaps — positioned as India's first AI-native hospital management platform.", size=12, bold=True, color=AI_PURPLE, font='Calibri')


def slide_05_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Our Solution: HMS Enterprise", "Cloud-native, AI-powered, regulation-ready")
    add_slide_footer(slide)

    # Hero description
    tb = add_text_box(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, "HMS Enterprise is a cloud-native, multi-tenant hospital management platform with 12 fully integrated modules, built on microservices architecture with Spring Boot, Apache Kafka, PostgreSQL, and deployed on AWS EKS.", size=13, color=BODY_TEXT, font='Calibri Light')

    # Key differentiators in cards
    diffs = [
        ("\U0001F3D7\uFE0F", "Schema-per-Tenant", "PostgreSQL schema isolation + RLS — zero cross-tenant data leakage", TEAL),
        ("\u26A1", "Event-Driven", "Apache Kafka backbone, 10K+ events/sec, real-time Socket.IO updates", AI_BLUE),
        ("\U0001F3E5", "ABDM-Native", "ABHA health ID, FHIR R4, bidirectional health record exchange from Day 1", GREEN),
        ("\U0001F916", "AI-First", "12 AI/ML modules — ambient AI, CDSS, revenue cycle AI, agentic workflows", AI_PURPLE),
        ("\U0001F6E1\uFE0F", "Regulation-Ready", "DPDP Act 2023, NABH 6th Ed, ABDM 2026 — compliance built into the core", GOLD),
        ("\U0001F4F1", "Mobile-First", "PWA patient portal, WhatsApp integration, offline-capable for Tier 2/3", NAVY),
    ]

    y = Inches(2.6)
    for i, (icon, title, desc, accent) in enumerate(diffs):
        col = Inches(0.5) + Inches(4.2) * (i % 3)
        if i == 3:
            y = Inches(4.5)
        card = add_shape(slide, col, y, Inches(3.9), Inches(1.5), RGBColor(247, 250, 252))
        add_shape(slide, col, y, Inches(3.9), Inches(0.05), accent)
        # Icon + Title
        tb_t = add_text_box(slide, col + Inches(0.15), y + Inches(0.15), Inches(3.5), Inches(0.4))
        set_text(tb_t.text_frame, f"{icon}  {title}", size=13, bold=True, color=NAVY)
        # Desc
        tb_d = add_text_box(slide, col + Inches(0.15), y + Inches(0.55), Inches(3.5), Inches(0.85))
        set_text(tb_d.text_frame, desc, size=10, color=BODY_TEXT, font='Calibri Light')


def slide_06_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Platform Architecture", "Cloud-native microservices with event-driven design")
    add_slide_footer(slide)

    # Architecture layers
    layers = [
        ("Presentation", "React 18 + TypeScript + PWA", TEAL),
        ("API Gateway", "Spring Cloud Gateway + Kong", TEAL),
        ("Business Logic", "Spring Boot 3.x (Java 21) — 12 Microservices", NAVY),
        ("Messaging", "Apache Kafka 3.x — Event Streaming", AI_PURPLE),
        ("Cache + Search", "Redis 7.x Cluster + Elasticsearch 8.x", AI_BLUE),
        ("Database", "PostgreSQL 16 — Schema-per-Tenant", DARK_NAVY),
        ("Infrastructure", "AWS EKS + S3 + CloudFront + WAF", DARK_NAVY),
    ]

    y = Inches(1.5)
    for label, tech, color in layers:
        # Layer bar
        add_shape(slide, Inches(0.7), y, Inches(7.5), Inches(0.55), color)
        tb = add_text_box(slide, Inches(0.9), y + Inches(0.05), Inches(2.5), Inches(0.45))
        set_text(tb.text_frame, label, size=11, bold=True, color=WHITE)
        tb2 = add_text_box(slide, Inches(3.3), y + Inches(0.05), Inches(4.7), Inches(0.45))
        set_text(tb2.text_frame, tech, size=10, color=RGBColor(220, 225, 235), font='Calibri Light')
        y += Inches(0.62)

    # Right side — Key metrics
    add_shape(slide, Inches(8.7), Inches(1.5), Inches(4.2), Inches(5.0), RGBColor(247, 250, 252))
    add_shape(slide, Inches(8.7), Inches(1.5), Inches(4.2), Inches(0.05), GOLD)

    tb_h = add_text_box(slide, Inches(8.9), Inches(1.7), Inches(3.8), Inches(0.4))
    set_text(tb_h.text_frame, "Architecture Highlights", size=14, bold=True, color=NAVY)

    metrics = [
        "70+ REST API endpoints",
        "6 Kafka topic namespaces",
        "10,000+ events/second",
        "< 500ms API P95 latency",
        "10K+ concurrent users/tenant",
        "50+ tenant capacity",
        "8-stage CI/CD pipeline",
        "Multi-AZ deployment",
        "99.9% uptime SLA",
        "4-hour RTO / 15-min RPO",
    ]
    tb_m = add_text_box(slide, Inches(8.9), Inches(2.2), Inches(3.8), Inches(4.0))
    tf = tb_m.text_frame
    tf.word_wrap = True
    set_text(tf, "", size=1)
    for m in metrics:
        add_bullet_para(tf, m, size=10, color=BODY_TEXT, font='Calibri Light', space_after=4)


def slide_07_modules_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "12 Integrated Modules", "Complete hospital workflow coverage from registration to discharge")
    add_slide_footer(slide)

    modules = [
        ("01", "OPD & Clinical", "\U0001F3E5", TEAL),
        ("02", "Pharmacy", "\U0001F48A", GREEN),
        ("03", "Laboratory", "\U0001F52C", AI_BLUE),
        ("04", "Finance & Billing", "\U0001F4B0", GOLD),
        ("05", "HR & Payroll", "\U0001F465", NAVY),
        ("06", "Inventory", "\U0001F4E6", ORANGE),
        ("07", "Business Intel", "\U0001F4CA", AI_PURPLE),
        ("08", "Patient Portal", "\U0001F4F1", TEAL),
        ("09", "Audit & Compliance", "\U0001F512", RED),
        ("10", "IPD Management", "\U0001F6CF\uFE0F", NAVY),
        ("11", "Scheduling", "\U0001F4C5", GREEN),
        ("12", "Platform Admin", "\u2699\uFE0F", DARK_NAVY),
    ]

    y_start = Inches(1.5)
    for i, (num, name, icon, color) in enumerate(modules):
        row = i // 4
        col_idx = i % 4
        x = Inches(0.5) + Inches(3.15) * col_idx
        y = y_start + Inches(1.7) * row
        # Card
        card = add_shape(slide, x, y, Inches(2.9), Inches(1.4), RGBColor(247, 250, 252))
        # Top accent
        add_shape(slide, x, y, Inches(2.9), Inches(0.05), color)
        # Number badge
        add_shape(slide, x + Inches(0.1), y + Inches(0.15), Inches(0.45), Inches(0.35), color)
        tb_n = add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(0.45), Inches(0.35))
        set_text(tb_n.text_frame, num, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Icon + Name
        tb_name = add_text_box(slide, x + Inches(0.65), y + Inches(0.15), Inches(2.1), Inches(0.4))
        set_text(tb_name.text_frame, f"{icon}  {name}", size=12, bold=True, color=NAVY)

    # Bottom callout
    callout = add_shape(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5), RGBColor(235, 248, 255))
    add_shape(slide, Inches(0.5), Inches(6.4), Inches(0.05), Inches(0.5), TEAL)
    tb_c = add_text_box(slide, Inches(0.8), Inches(6.42), Inches(11.7), Inches(0.45))
    set_text(tb_c.text_frame, "All modules communicate via Apache Kafka events — enabling real-time cross-module workflows and unified patient records.", size=11, italic=True, color=TEAL, font='Calibri Light')


def slide_08_modules_clinical(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Clinical Modules", "OPD, Pharmacy, Laboratory, IPD")
    add_slide_footer(slide)

    mods = [
        ("OPD & Clinical Operations", [
            "SOAP-based EMR with clinical templates",
            "Token queue with real-time wait times",
            "E-prescriptions with drug interaction checks",
            "Vitals trending with alert thresholds",
            "Referral management & follow-ups",
        ], TEAL),
        ("Pharmacy Management", [
            "FIFO batch tracking & expiry alerts",
            "Drug interaction engine (real-time)",
            "GST-compliant billing (HSN codes)",
            "Controlled substance (Schedule H/H1)",
            "Auto-receive prescriptions from EMR",
        ], GREEN),
        ("Laboratory (LIS)", [
            "Barcode-tracked sample lifecycle",
            "Dual validation (tech + pathologist)",
            "Auto-analyzer integration (50+ models)",
            "Critical value flagging & auto-notify",
            "NABL-aligned workflows",
        ], AI_BLUE),
        ("IPD Management", [
            "Real-time bed board & housekeeping",
            "Nursing records & care plans",
            "MAR with missed-dose alerts",
            "Clinical pathway tracking",
            "Checklist-based discharge planning",
        ], NAVY),
    ]

    for i, (title, features, color) in enumerate(mods):
        col = Inches(0.4) + Inches(3.2) * i
        # Card bg
        add_shape(slide, col, Inches(1.5), Inches(3.0), Inches(5.2), RGBColor(247, 250, 252))
        add_shape(slide, col, Inches(1.5), Inches(3.0), Inches(0.05), color)
        # Title
        tb = add_text_box(slide, col + Inches(0.1), Inches(1.65), Inches(2.8), Inches(0.5))
        set_text(tb.text_frame, title, size=11, bold=True, color=NAVY)
        # Features
        tb_f = add_text_box(slide, col + Inches(0.1), Inches(2.2), Inches(2.8), Inches(4.2))
        tf = tb_f.text_frame
        tf.word_wrap = True
        set_text(tf, "", size=1)
        for feat in features:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = "\u2713  " + feat
            run.font.size = Pt(9.5)
            run.font.color.rgb = BODY_TEXT
            run.font.name = 'Calibri Light'
            p.space_after = Pt(6)


def slide_09_modules_ops(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Operations & Intelligence", "Finance, HR, Inventory, BI, Patient Portal, Scheduling, Audit, Admin")
    add_slide_footer(slide)

    mods = [
        ("Finance & Billing", "Multi-module bill aggregation, TPA/insurance workflows, payment gateway (Razorpay/UPI), GST compliance, revenue dashboards.", GOLD),
        ("HR & Payroll", "Biometric attendance, shift scheduling, PF/ESIC/TDS deductions, leave management, employee lifecycle.", NAVY),
        ("Inventory & Procurement", "Auto-reorder, vendor management, PO workflows, GRN processing, equipment AMC tracking.", ORANGE),
        ("Business Intelligence", "50+ KPIs, Elasticsearch dashboards, trend analysis, custom report builder, NABH benchmarks.", AI_PURPLE),
        ("Patient Portal (PWA)", "Self-service booking, health records access, digital payments, notifications, NPS feedback.", TEAL),
        ("Scheduling", "Multi-channel booking (web, WhatsApp, IVR), no-show prediction, queue optimization, OT booking.", GREEN),
        ("Audit & Compliance", "Immutable audit trail, DPDP consent management, NABH tracking, 7-year retention, security monitoring.", RED),
        ("Platform Admin", "Tenant onboarding, feature flags, cross-tenant analytics, system health, billing & subscriptions.", DARK_NAVY),
    ]

    y = Inches(1.5)
    for i, (title, desc, color) in enumerate(mods):
        col = Inches(0.4) if i % 2 == 0 else Inches(6.8)
        if i > 0 and i % 2 == 0:
            y += Inches(1.2)
        card = add_shape(slide, col, y, Inches(6.1), Inches(1.05), RGBColor(247, 250, 252))
        add_shape(slide, col, y, Inches(0.05), Inches(1.05), color)
        tb = add_text_box(slide, col + Inches(0.2), y + Inches(0.05), Inches(5.6), Inches(0.3))
        set_text(tb.text_frame, title, size=12, bold=True, color=NAVY)
        tb2 = add_text_box(slide, col + Inches(0.2), y + Inches(0.4), Inches(5.6), Inches(0.6))
        set_text(tb2.text_frame, desc, size=9.5, color=BODY_TEXT, font='Calibri Light')
    y += Inches(0.0)


def slide_10_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Security & Compliance", "Enterprise-grade protection built into every layer")
    add_slide_footer(slide)

    # RBAC section
    add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.4), DARK_NAVY)
    tb_h = add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4))
    set_text(tb_h.text_frame, "6-Tier RBAC Hierarchy", size=12, bold=True, color=WHITE)

    roles = [
        ("L0 SuperAdmin", "Platform-wide", DARK_NAVY),
        ("L1 TenantAdmin", "Single tenant", NAVY),
        ("L2 Branch Lead", "Single branch", TEAL),
        ("L3 Dept Head", "Department", AI_BLUE),
        ("L4 Staff", "Assigned resources", GREEN),
        ("L5 Patient", "Own records", GRAY),
    ]
    y = Inches(2.05)
    for role, scope, color in roles:
        add_shape(slide, Inches(0.5), y, Inches(0.08), Inches(0.35), color)
        tb = add_text_box(slide, Inches(0.7), y, Inches(2.5), Inches(0.35))
        set_text(tb.text_frame, role, size=9.5, bold=True, color=NAVY, font='Calibri')
        tb2 = add_text_box(slide, Inches(3.3), y, Inches(3.0), Inches(0.35))
        set_text(tb2.text_frame, scope, size=9.5, color=GRAY, font='Calibri Light')
        y += Inches(0.38)

    # Compliance section
    add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4), DARK_NAVY)
    tb_h2 = add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5.6), Inches(0.4))
    set_text(tb_h2.text_frame, "Regulatory Compliance", size=12, bold=True, color=WHITE)

    compliances = [
        ("\u2705", "DPDP Act 2023", "Consent management, data minimization, right to erasure, breach notification"),
        ("\u2705", "ABDM 2026", "ABHA Health ID, FHIR R4, Health Information Exchange, UHI APIs"),
        ("\u2705", "NABH 6th Edition", "Quality indicators, gap analysis, evidence management, compliance scoring"),
        ("\u2705", "ISO 27001:2022", "Information security management system"),
        ("\u2705", "OWASP Top 10", "Web application security — JWT, MFA, HttpOnly cookies, RLS"),
    ]
    y2 = Inches(2.05)
    for icon, title, desc in compliances:
        tb = add_text_box(slide, Inches(7.1), y2, Inches(5.5), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = f"{icon}  {title}"
        run1.font.size = Pt(10)
        run1.font.bold = True
        run1.font.color.rgb = NAVY
        run1.font.name = 'Calibri'
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(8.5)
        run2.font.color.rgb = GRAY
        run2.font.name = 'Calibri Light'
        p2.space_after = Pt(2)
        y2 += Inches(0.72)

    # Security callout
    callout = add_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.9), RGBColor(255, 245, 245))
    add_shape(slide, Inches(0.5), Inches(5.5), Inches(0.06), Inches(0.9), RED)
    tb_c = add_text_box(slide, Inches(0.8), Inches(5.55), Inches(11.8), Inches(0.8))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    set_text(tf_c, "\U0001F510  Data Protection: AES-256-GCM field-level encryption for all PII. Tenant-specific keys rotated every 90 days. All AI/ML calls use PII de-identification — patient data stripped before external API calls.", size=10.5, color=RGBColor(130, 30, 30), font='Calibri')


def slide_11_ai_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "AI/ML Innovation Roadmap", "12 AI modules across 4 phases — from predictive analytics to agentic workflows")
    add_slide_footer(slide)

    # Stat boxes
    stats = [("12", "AI Modules"), ("4", "Delivery Phases"), ("18", "Gaps Addressed"), ("\u20B917,500Cr", "AI Health Mkt 2030")]
    x = Inches(0.5)
    for val, label in stats:
        add_stat_box(slide, x, Inches(1.4), val, label, Inches(2.95))
        x += Inches(3.15)

    phases = [
        ("Phase 1: Foundation (M1\u20136)", [
            "Predictive Bed Occupancy Engine",
            "Intelligent Appointment Optimizer",
            "Smart Inventory Forecasting",
        ], TEAL),
        ("Phase 2: Clinical (M7\u201312)", [
            "Clinical Decision Support (CDSS)",
            "Ambient Clinical Documentation",
            "Diagnostic Imaging AI",
        ], AI_BLUE),
        ("Phase 3: Revenue (M13\u201318)", [
            "AI Revenue Cycle Management",
            "Patient Risk Stratification",
            "Discharge Summary Generator",
        ], AI_PURPLE),
        ("Phase 4: Agentic (M19\u201324)", [
            "Agentic Clinical Workflows",
            "Conversational Patient AI",
            "Federated Learning Network",
        ], DARK_NAVY),
    ]

    x = Inches(0.4)
    for title, items, color in phases:
        card = add_shape(slide, x, Inches(2.9), Inches(3.0), Inches(3.6), RGBColor(247, 250, 252))
        add_shape(slide, x, Inches(2.9), Inches(3.0), Inches(0.06), color)
        # Phase title
        add_shape(slide, x, Inches(2.96), Inches(3.0), Inches(0.4), color)
        tb = add_text_box(slide, x + Inches(0.1), Inches(2.98), Inches(2.8), Inches(0.38))
        set_text(tb.text_frame, title, size=10, bold=True, color=WHITE)
        # Items
        tb_i = add_text_box(slide, x + Inches(0.15), Inches(3.5), Inches(2.7), Inches(2.8))
        tf = tb_i.text_frame
        tf.word_wrap = True
        set_text(tf, "", size=1)
        for item in items:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = "\u25B8  " + item
            run.font.size = Pt(10)
            run.font.color.rgb = NAVY
            run.font.name = 'Calibri'
            p.space_after = Pt(10)
        x += Inches(3.2)


def slide_12_pricing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Pricing Model", "Tiered SaaS pricing designed for every hospital size")
    add_slide_footer(slide)

    tiers = [
        ("STARTER", "\u20B982K\u2013\u20B92.1L/mo", "20\u2013100 beds", [
            "6 core modules included",
            "Email support (48h SLA)",
            "500 req/min API limit",
            "Self-service onboarding",
            "99.5% SLA",
        ], TEAL, RGBColor(230, 255, 250)),
        ("PROFESSIONAL", "\u20B94.2L\u2013\u20B912.5L/mo", "100\u2013500 beds", [
            "All 12 modules included",
            "Basic AI features",
            "Priority support (4h SLA)",
            "2,000 req/min API limit",
            "99.9% SLA",
        ], AI_BLUE, RGBColor(235, 240, 255)),
        ("ENTERPRISE", "\u20B912.5L\u2013\u20B921L/mo", "500+ beds", [
            "Full AI/ML suite",
            "Dedicated CSM (1h SLA)",
            "5,000 req/min API limit",
            "White-glove onboarding",
            "99.95% SLA",
        ], GOLD, RGBColor(255, 255, 240)),
    ]

    x = Inches(0.7)
    for name, price, beds, features, color, bg in tiers:
        w = Inches(3.8)
        # Card
        card = add_shape(slide, x, Inches(1.5), w, Inches(4.8), bg)
        # Top accent
        add_shape(slide, x, Inches(1.5), w, Inches(0.06), color)
        # Tier name
        add_shape(slide, x, Inches(1.56), w, Inches(0.45), color)
        tb_n = add_text_box(slide, x, Inches(1.58), w, Inches(0.42))
        set_text(tb_n.text_frame, name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Price
        tb_p = add_text_box(slide, x, Inches(2.15), w, Inches(0.5))
        set_text(tb_p.text_frame, price, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Beds
        tb_b = add_text_box(slide, x, Inches(2.6), w, Inches(0.3))
        set_text(tb_b.text_frame, beds, size=11, color=GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)
        # Divider
        add_shape(slide, x + Inches(0.3), Inches(3.0), w - Inches(0.6), Inches(0.02), LIGHT_GRAY)
        # Features
        tb_f = add_text_box(slide, x + Inches(0.3), Inches(3.15), w - Inches(0.6), Inches(3.0))
        tf = tb_f.text_frame
        tf.word_wrap = True
        set_text(tf, "", size=1)
        for feat in features:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = "\u2713  " + feat
            run.font.size = Pt(10)
            run.font.color.rgb = BODY_TEXT
            run.font.name = 'Calibri Light'
            p.space_after = Pt(8)
        x += Inches(4.1)


def slide_13_projections(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Financial Projections", "5-year growth trajectory")
    add_slide_footer(slide)

    # Year columns as visual bars
    years = [
        ("Year 1", "25", "\u20B912.5 Cr", "65%", TEAL, 0.5),
        ("Year 2", "100", "\u20B967 Cr", "72%", AI_BLUE, 1.2),
        ("Year 3", "350", "\u20B9234 Cr", "78%", AI_PURPLE, 2.5),
        ("Year 5", "1,200", "\u20B91,000 Cr", "82%", GOLD, 5.0),
    ]

    # Bar chart visual
    x = Inches(1.0)
    max_h = Inches(3.5)
    for yr, hospitals, arr, margin, color, scale in years:
        bar_h = Inches(scale * 0.7)
        bar_y = Inches(5.5) - bar_h
        # Bar
        add_shape(slide, x, bar_y, Inches(2.2), bar_h, color)
        # Value on bar
        tb_v = add_text_box(slide, x, bar_y + Inches(0.05), Inches(2.2), Inches(0.35))
        set_text(tb_v.text_frame, arr, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Year label below
        tb_yr = add_text_box(slide, x, Inches(5.55), Inches(2.2), Inches(0.35))
        set_text(tb_yr.text_frame, yr, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # Hospitals
        tb_h = add_text_box(slide, x, Inches(5.9), Inches(2.2), Inches(0.25))
        set_text(tb_h.text_frame, f"{hospitals} hospitals", size=9, color=GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)
        # Margin
        tb_m = add_text_box(slide, x, Inches(6.15), Inches(2.2), Inches(0.25))
        set_text(tb_m.text_frame, f"{margin} gross margin", size=9, color=GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)
        x += Inches(2.8)

    # Revenue streams
    tb_rs = add_text_box(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.4))
    set_text(tb_rs.text_frame, "Revenue Streams: Subscription (70%) + AI Add-ons (15%) + Implementation (10%) + Marketplace (5%)", size=11, color=BODY_TEXT, font='Calibri Light')


def slide_14_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, "Implementation Roadmap", "54-week phased delivery with continuous value")
    add_slide_footer(slide)

    phases = [
        ("Phase 1", "Weeks 1\u201314", "Core Platform", "OPD, Pharmacy, Lab, Finance, IPD, Scheduling\nMVP launch with 3 pilot hospitals", TEAL),
        ("Phase 2", "Weeks 15\u201328", "Extended Modules", "HR, Inventory, BI, Patient Portal, Admin\nFeature-complete platform release", AI_BLUE),
        ("Phase 3", "Weeks 29\u201342", "AI & Intelligence", "CDSS, Ambient AI, Predictive Analytics\nAI-powered clinical features live", AI_PURPLE),
        ("Phase 4", "Weeks 43\u201354", "Scale & Optimize", "RCM AI, Agentic Workflows, Marketplace\nFull platform with 50+ tenants", GOLD),
    ]

    # Timeline bar
    add_shape(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.08), LIGHT_GRAY)

    x = Inches(0.5)
    for phase, weeks, title, desc, color in phases:
        w = Inches(3.0)
        # Phase segment on timeline
        add_shape(slide, x, Inches(3.4), w, Inches(0.28), color)
        # Phase label
        tb_p = add_text_box(slide, x, Inches(3.42), w, Inches(0.25))
        set_text(tb_p.text_frame, f"{phase}: {weeks}", size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Circle connector
        add_shape(slide, x + Inches(1.35), Inches(3.75), Inches(0.3), Inches(0.3), color)
        # Card above or below
        if phases.index((phase, weeks, title, desc, color)) % 2 == 0:
            card_y = Inches(1.5)
        else:
            card_y = Inches(4.3)
        card = add_shape(slide, x, card_y, w, Inches(1.7), RGBColor(247, 250, 252))
        add_shape(slide, x, card_y, w, Inches(0.05), color)
        tb_t = add_text_box(slide, x + Inches(0.1), card_y + Inches(0.15), w - Inches(0.2), Inches(0.35))
        set_text(tb_t.text_frame, title, size=13, bold=True, color=NAVY)
        tb_d = add_text_box(slide, x + Inches(0.1), card_y + Inches(0.55), w - Inches(0.2), Inches(1.0))
        set_text(tb_d.text_frame, desc, size=9.5, color=BODY_TEXT, font='Calibri Light', line_spacing=1.5)
        x += Inches(3.15)

    # Team summary
    tb_team = add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.35))
    set_text(tb_team.text_frame, "Team: 1 Product Owner  \u2022  1 Architect  \u2022  4 Backend  \u2022  3 Frontend  \u2022  2 AI/ML  \u2022  2 QA  \u2022  1 DevOps  \u2022  1 UI/UX  \u2022  1 PM  =  16 Members", size=10, color=NAVY, font='Calibri', align=PP_ALIGN.CENTER)


def slide_15_why_us(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_slide_header(slide, f"Why {COMPANY_FULL}?", "The right partner for India's healthcare transformation")
    add_slide_footer(slide)

    reasons = [
        ("\U0001F947", "First-Mover in AI Healthcare", "Only Indian-origin platform offering ambient clinical AI, CDSS, and AI-powered revenue cycle management. 18 market gaps addressed."),
        ("\U0001F3D7\uFE0F", "True Multi-Tenant SaaS", "Schema-per-tenant isolation with 3-layer defense. Not a hosted single-tenant system pretending to be SaaS."),
        ("\U0001F512", "Regulation-Ready from Day 1", "DPDP Act, ABDM, NABH built into the core — not bolted on as afterthoughts. Future-proof for 2026 mandates."),
        ("\u26A1", "Modern Tech Stack", "Spring Boot + Kafka + PostgreSQL + Redis + EKS. Event-driven, cloud-native, auto-scaling. Not legacy monolith."),
        ("\U0001F4CA", "Data-Driven Operations", "Elasticsearch BI, 50+ KPIs, predictive analytics. Turn hospital data into actionable insights."),
        ("\U0001F91D", "Partnership Model", "Not just software — white-glove onboarding, dedicated CSM, 12-week warranty, 24/7 critical support."),
    ]

    y = Inches(1.5)
    for i, (icon, title, desc) in enumerate(reasons):
        col = Inches(0.4) if i % 2 == 0 else Inches(6.8)
        if i > 0 and i % 2 == 0:
            y += Inches(1.55)
        card = add_shape(slide, col, y, Inches(6.1), Inches(1.35), RGBColor(247, 250, 252))
        add_shape(slide, col, y, Inches(0.06), Inches(1.35), TEAL if i % 3 == 0 else AI_PURPLE if i % 3 == 1 else GOLD)
        tb_t = add_text_box(slide, col + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.35))
        set_text(tb_t.text_frame, f"{icon}  {title}", size=13, bold=True, color=NAVY)
        tb_d = add_text_box(slide, col + Inches(0.2), y + Inches(0.5), Inches(5.6), Inches(0.8))
        set_text(tb_d.text_frame, desc, size=10, color=BODY_TEXT, font='Calibri Light')


def slide_16_thank_you(prs):
    """Closing slide with branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK_NAVY)

    # Top teal accent
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), TEAL)

    # Company
    tb = add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.5))
    set_text(tb.text_frame, COMPANY_FULL.upper(), size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Gold divider
    add_shape(slide, Inches(4.5), Inches(2.2), Inches(4.3), Inches(0.04), GOLD)

    # Thank You
    tb2 = add_text_box(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(1.0))
    set_text(tb2.text_frame, "Thank You", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    tb3 = add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6))
    set_text(tb3.text_frame, "Let's Transform Healthcare Together", size=20, color=LIGHT_PURPLE, font='Calibri Light', align=PP_ALIGN.CENTER)

    # Gold divider
    add_shape(slide, Inches(4.5), Inches(4.4), Inches(4.3), Inches(0.04), GOLD)

    # Contact info
    tb4 = add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.4))
    set_text(tb4.text_frame, f"Author: {AUTHOR}", size=13, color=WHITE, font='Calibri Light', align=PP_ALIGN.CENTER)

    tb5 = add_text_box(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.4))
    set_text(tb5.text_frame, f"{COMPANY_FULL}  |  {PRODUCT_NAME}", size=11, color=SOFT_GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)

    tb6 = add_text_box(slide, Inches(2), Inches(5.6), Inches(9.3), Inches(0.3))
    set_text(tb6.text_frame, DOC_DATE, size=10, color=GRAY, font='Calibri Light', align=PP_ALIGN.CENTER)

    # Confidential
    tb7 = add_text_box(slide, Inches(2), Inches(6.3), Inches(9.3), Inches(0.3))
    set_text(tb7.text_frame, "CONFIDENTIAL", size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Bottom teal
    add_shape(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), TEAL)


# ═══════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()

    # Set 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Building slide 01: Title...")
    slide_01_title(prs)

    print("Building slide 02: Agenda...")
    slide_02_agenda(prs)

    print("Building slide 03: The Problem...")
    slide_03_problem(prs)

    print("Building slide 04: Market Opportunity...")
    slide_04_market(prs)

    print("Building slide 05: Our Solution...")
    slide_05_solution(prs)

    print("Building slide 06: Platform Architecture...")
    slide_06_architecture(prs)

    print("Building slide 07: 12 Modules Overview...")
    slide_07_modules_overview(prs)

    print("Building slide 08: Clinical Modules...")
    slide_08_modules_clinical(prs)

    print("Building slide 09: Operations Modules...")
    slide_09_modules_ops(prs)

    print("Building slide 10: Security & Compliance...")
    slide_10_security(prs)

    print("Building slide 11: AI/ML Roadmap...")
    slide_11_ai_roadmap(prs)

    print("Building slide 12: Pricing...")
    slide_12_pricing(prs)

    print("Building slide 13: Financial Projections...")
    slide_13_projections(prs)

    print("Building slide 14: Implementation Timeline...")
    slide_14_timeline(prs)

    print("Building slide 15: Why TechDigital WishTree...")
    slide_15_why_us(prs)

    print("Building slide 16: Thank You...")
    slide_16_thank_you(prs)

    prs.save(OUTPUT_PATH)
    print(f"\nPresentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")
    print("Done!")


if __name__ == "__main__":
    main()
